from io import StringIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from random import choice
import random
import sys


class OutputInterceptor(list):
    def __enter__(self):
        self._stdout = sys.stdout
        sys.stdout = self._stringio = StringIO()
        return self

    def __exit__(self, *args):
        self.extend(self._stringio.getvalue().splitlines())
        del self._stringio
        sys.stdout = self._stdout


prs = Presentation()
for i in range(5):
    title_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    left = Inches(1) 
    top = Inches(1)
    width = Inches(8)
    height = Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    method = choice(dir(random))
    title.text = method
    with OutputInterceptor() as description:
        help('random.' + method)
    # subtitle.text = '\n'.join(description)
    tf = txBox.text_frame
    for desc in description:
        p = tf.add_paragraph()
        p.text = desc
        p.font.name = 'Courier New'
        p.font.size = Pt(14)
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    # subtitle = slide.placeholders[1]
prs.save('res.pptx')
# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#add-textbox-example