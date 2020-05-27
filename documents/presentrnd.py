import random
from pptx import Presentation

prs = Presentation()

for method in dir(random):
    print(help('random.' + method))
    sl = prs.Slides.add_slide()
    sl.shapes.title = method
    

